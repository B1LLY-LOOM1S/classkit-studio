import streamlit as st
import sqlite3
import uuid
import os
import json
import io
import datetime
import textwrap
from dotenv import load_dotenv

# Document Generation Libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Pt as DocxPt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Frame, PageTemplate
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch

# AI
import google.generativeai as genai

# Load Env
load_dotenv()

# --- CONFIGURATION ---
DB_FILE = "classkit.db"
TEACHER_CODE = os.getenv("TEACHER_CODE", "changeme")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")

st.set_page_config(page_title="ClassKit Studio", layout="wide", page_icon="🎓")

# --- DATABASE LAYER ---
def init_db():
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    c.execute('''
        CREATE TABLE IF NOT EXISTS projects (
            id TEXT PRIMARY KEY,
            created_at TEXT,
            title TEXT,
            subject TEXT,
            grade TEXT,
            source_notes TEXT,
            slides_json TEXT,
            poster_json TEXT,
            assignment_json TEXT,
            teacher_token TEXT,
            student_token TEXT
        )
    ''')
    conn.commit()
    conn.close()

def save_project(data):
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    # Check if exists
    c.execute("SELECT id FROM projects WHERE id = ?", (data['id'],))
    exists = c.fetchone()
    
    if exists:
        c.execute('''
            UPDATE projects SET 
                title=?, subject=?, grade=?, source_notes=?, 
                slides_json=?, poster_json=?, assignment_json=?
            WHERE id=?
        ''', (
            data['title'], data['subject'], data['grade'], data['source_notes'],
            json.dumps(data.get('slides_json', {})),
            json.dumps(data.get('poster_json', {})),
            json.dumps(data.get('assignment_json', {})),
            data['id']
        ))
    else:
        c.execute('''
            INSERT INTO projects VALUES (?,?,?,?,?,?,?,?,?,?,?)
        ''', (
            data['id'], datetime.datetime.now().isoformat(),
            data['title'], data['subject'], data['grade'], data['source_notes'],
            json.dumps(data.get('slides_json', {})),
            json.dumps(data.get('poster_json', {})),
            json.dumps(data.get('assignment_json', {})),
            data['teacher_token'], data['student_token']
        ))
    conn.commit()
    conn.close()

def load_project(project_id):
    conn = sqlite3.connect(DB_FILE)
    conn.row_factory = sqlite3.Row
    c = conn.cursor()
    c.execute("SELECT * FROM projects WHERE id = ?", (project_id,))
    row = c.fetchone()
    conn.close()
    if row:
        d = dict(row)
        # Parse JSON fields back to dicts
        d['slides_json'] = json.loads(d['slides_json']) if d['slides_json'] else {}
        d['poster_json'] = json.loads(d['poster_json']) if d['poster_json'] else {}
        d['assignment_json'] = json.loads(d['assignment_json']) if d['assignment_json'] else {}
        return d
    return None

def get_project_by_token(token, mode="student"):
    conn = sqlite3.connect(DB_FILE)
    conn.row_factory = sqlite3.Row
    c = conn.cursor()
    col = "student_token" if mode == "student" else "teacher_token"
    c.execute(f"SELECT * FROM projects WHERE {col} = ?", (token,))
    row = c.fetchone()
    conn.close()
    if row:
        return load_project(row['id'])
    return None

# --- AI LAYER ---
def get_mock_data(prompt_type):
    """Fallback if no API key"""
    if prompt_type == "slides":
        return {
            "deck_title": "Demo: The Solar System",
            "slides": [
                {"type": "title", "title": "The Solar System", "bullets": [], "speaker_notes": "Intro"},
                {"type": "content", "title": "Inner Planets", "bullets": ["Mercury", "Venus", "Earth", "Mars"], "speaker_notes": "Rocky planets."},
                {"type": "summary", "title": "Review", "bullets": ["8 Planets", "Sun is a star"], "speaker_notes": "Wrap up."}
            ]
        }
    elif prompt_type == "poster":
        return {
            "poster_title": "Lab Safety Rules",
            "sections": [
                {"heading": "Protection", "body_bullets": ["Wear Goggles", "Use Gloves"]},
                {"heading": "Behavior", "body_bullets": ["No running", "No eating"]}
            ],
            "footer_callout": "Stay Safe!"
        }
    elif prompt_type == "assignment":
        return {
            "assignment_title": "Solar System Quiz",
            "instructions": "Answer all questions.",
            "questions": [
                {"type": "mcq", "prompt": "Which is the red planet?", "choices": ["Mars", "Venus"], "answer": "Mars", "explanation": "Iron oxide dust."},
                {"type": "short", "prompt": "Name the largest planet.", "choices": [], "answer": "Jupiter", "explanation": "Gas giant."}
            ],
            "rubric": ["1pt per correct answer"]
        }
    return {}

def call_gemini_json(prompt, schema_hint):
    if not GOOGLE_API_KEY:
        # Return mock data based on simple keyword matching for robustness
        if "slides" in prompt.lower(): return get_mock_data("slides")
        if "poster" in prompt.lower(): return get_mock_data("poster")
        return get_mock_data("assignment")

    genai.configure(api_key=GOOGLE_API_KEY)
    model = genai.GenerativeModel('gemini-pro')
    
    full_prompt = f"""
    You are an educational content generator.
    Output strictly valid JSON. No markdown fences.
    
    Task: {prompt}
    
    Required JSON Structure:
    {schema_hint}
    """
    
    try:
        response = model.generate_content(full_prompt)
        txt = response.text.strip()
        # Cleanup markdown if present
        if txt.startswith("```"):
            lines = txt.splitlines()
            if lines[0].startswith("```"): lines = lines[1:]
            if lines[-1].startswith("```"): lines = lines[:-1]
            txt = "\n".join(lines)
        return json.loads(txt)
    except Exception as e:
        st.error(f"AI Generation failed: {e}")
        return {}

# --- EXPORT HANDLERS ---
def generate_pptx_file(data):
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = data.get('deck_title', 'Untitled')
    slide.placeholders[1].text = "Generated by ClassKit Studio"

    for s_data in data.get('slides', []):
        if s_data.get('type') == 'title':
            layout = prs.slide_layouts[0]
        else:
            layout = prs.slide_layouts[1] # Title and Content
            
        slide = prs.slides.add_slide(layout)
        
        # Set Title
        if slide.shapes.title:
            slide.shapes.title.text = s_data.get('title', '')
        
        # Set Content
        if s_data.get('type') != 'title' and 'bullets' in s_data:
            # Find body placeholder
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    tf = shape.text_frame
                    tf.clear()
                    for bullet in s_data['bullets']:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
        
        # Notes
        if 'speaker_notes' in s_data:
            slide.notes_slide.notes_text_frame.text = s_data['speaker_notes']

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

def generate_pdf_poster(data):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=30, leftMargin=30, topMargin=30, bottomMargin=30)
    styles = getSampleStyleSheet()
    
    # Custom Styles
    title_style = ParagraphStyle('PosterTitle', parent=styles['Heading1'], fontSize=28, alignment=1, textColor=colors.darkblue, spaceAfter=20)
    head_style = ParagraphStyle('SectionHead', parent=styles['Heading2'], fontSize=16, textColor=colors.white, backColor=colors.steelblue, borderPadding=5, spaceAfter=10)
    body_style = ParagraphStyle('Body', parent=styles['Normal'], fontSize=12, leading=16)
    footer_style = ParagraphStyle('Footer', parent=styles['Italic'], alignment=1, fontSize=10, textColor=colors.grey)

    story = []
    
    # Header
    story.append(Paragraph(data.get('poster_title', 'Class Poster'), title_style))
    story.append(Spacer(1, 0.2*inch))
    
    # Sections
    sections = data.get('sections', [])
    for sec in sections:
        # Section Header Box
        story.append(Paragraph(sec.get('heading', ''), head_style))
        
        # Bullets
        bullets = sec.get('body_bullets', [])
        for b in bullets:
            story.append(Paragraph(f"• {b}", body_style))
        story.append(Spacer(1, 0.2*inch))
        
    # Footer
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph(data.get('footer_callout', ''), footer_style))
    story.append(Paragraph("Generated by ClassKit Studio - Educational Use Only", footer_style))

    doc.build(story)
    return buffer.getvalue()

def generate_docx_assignment(data, include_answers=False):
    doc = Document()
    
    # Header
    title = doc.add_heading(data.get('assignment_title', 'Assignment'), 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    if include_answers:
        p = doc.add_paragraph()
        run = p.add_run("TEACHER COPY - ANSWER KEY")
        run.bold = True
        run.font.color.rgb = RGBColor(255, 0, 0)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph(data.get('instructions', ''))
    doc.add_paragraph("_" * 50)
    
    # Questions
    for idx, q in enumerate(data.get('questions', []), 1):
        doc.add_heading(f"Q{idx}: {q.get('prompt')}", level=2)
        
        if q.get('type') == 'mcq':
            for choice in q.get('choices', []):
                doc.add_paragraph(f"[ ] {choice}", style='List Bullet')
        elif q.get('type') == 'short':
            doc.add_paragraph("\n\n__________________________\n")
        
        if include_answers:
            ans_para = doc.add_paragraph()
            ans_para.add_run(f"\nAnswer: {q.get('answer')}").bold = True
            ans_para.add_run(f"\nExplanation: {q.get('explanation')}")
            
    # Rubric (Only in teacher key usually, but requested in output)
    if include_answers and 'rubric' in data:
        doc.add_page_break()
        doc.add_heading("Rubric", level=1)
        for r in data['rubric']:
            doc.add_paragraph(r, style='List Bullet')

    # Footer/Watermark
    section = doc.sections[0]
    footer = section.footer
    p = footer.paragraphs[0]
    p.text = "Generated by ClassKit Studio - Educational Use Only"
    if include_answers:
        p.text += " | TEACHER COPY"

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()

# --- APP UI ---
def render_teacher_view(project):
    st.subheader(f"Project: {project['title']}")
    
    # Main Tabs
    tab1, tab2, tab3 = st.tabs(["📊 Slides", "🖼️ Poster", "📝 Assignment"])
    
    with tab1:
        st.markdown("### Slide Deck Configuration")
        if st.button("Generate Slides (AI)"):
            with st.spinner("Generating slides..."):
                schema = """
                {
                    "deck_title": "string",
                    "slides": [
                        {"type": "title|content|summary", "title": "string", "bullets": ["string"], "speaker_notes": "string"}
                    ]
                }
                """
                prompt = f"Create a slide deck for {project['grade']} grade {project['subject']} about: {project['source_notes']}"
                res = call_gemini_json(prompt, schema)
                if res:
                    project['slides_json'] = res
                    save_project(project)
                    st.success("Slides Generated!")
                    st.rerun()

        if project['slides_json']:
            st.json(project['slides_json'], expanded=False)
            pptx_bytes = generate_pptx_file(project['slides_json'])
            st.download_button("Download .pptx", pptx_bytes, "slides.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        else:
            st.info("No slides generated yet.")

    with tab2:
        st.markdown("### Class Poster")
        if st.button("Generate Poster (AI)"):
            with st.spinner("Designing poster..."):
                schema = """
                {
                    "poster_title": "string",
                    "sections": [{"heading": "string", "body_bullets": ["string"]}],
                    "footer_callout": "string"
                }
                """
                prompt = f"Create a one-page educational poster content for {project['grade']} grade {project['subject']} about: {project['source_notes']}"
                res = call_gemini_json(prompt, schema)
                if res:
                    project['poster_json'] = res
                    save_project(project)
                    st.success("Poster Designed!")
                    st.rerun()
        
        if project['poster_json']:
            st.json(project['poster_json'], expanded=False)
            pdf_bytes = generate_pdf_poster(project['poster_json'])
            st.download_button("Download Poster PDF", pdf_bytes, "poster.pdf", "application/pdf")
        else:
            st.info("No poster generated yet.")

    with tab3:
        st.markdown("### Assignment & Answer Key")
        if st.button("Generate Assignment (AI)"):
            with st.spinner("Writing questions..."):
                schema = """
                {
                    "assignment_title": "string",
                    "instructions": "string",
                    "questions": [{"type": "mcq|short", "prompt": "string", "choices": ["string"], "answer": "string", "explanation": "string"}],
                    "rubric": ["string"]
                }
                """
                prompt = f"Create a homework assignment with 3 MCQs and 2 Short Answers for {project['grade']} grade {project['subject']} about: {project['source_notes']}"
                res = call_gemini_json(prompt, schema)
                if res:
                    project['assignment_json'] = res
                    save_project(project)
                    st.success("Assignment Created!")
                    st.rerun()

        if project['assignment_json']:
            st.json(project['assignment_json'], expanded=False)
            col1, col2 = st.columns(2)
            with col1:
                docx_student = generate_docx_assignment(project['assignment_json'], include_answers=False)
                st.download_button("Download Student Doc (.docx)", docx_student, "assignment_student.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
            with col2:
                docx_key = generate_docx_assignment(project['assignment_json'], include_answers=True)
                st.download_button("Download Answer Key (.docx)", docx_key, "assignment_key.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        else:
            st.info("No assignment generated yet.")
            
    st.markdown("---")
    st.markdown("#### Share with Students")
    link = f"http://localhost:8501/?token={project['student_token']}"
    st.code(link, language="text")

def render_student_view(project):
    st.title(f"📚 {project['title']}")
    st.info(f"Subject: {project['subject']} | Grade: {project['grade']}")
    
    tab_assign, tab_poster = st.tabs(["Assignment", "Reference Poster"])
    
    with tab_assign:
        if project['assignment_json']:
            data = project['assignment_json']
            st.subheader(data.get('assignment_title'))
            st.write(data.get('instructions'))
            for idx, q in enumerate(data.get('questions', []), 1):
                st.markdown(f"**Q{idx}: {q.get('prompt')}**")
                if q.get('type') == 'mcq':
                    st.radio(f"Select answer for Q{idx}", q.get('choices', []), key=f"q_{idx}")
                else:
                    st.text_area(f"Your answer for Q{idx}", key=f"q_{idx}")
                st.divider()
            
            docx_student = generate_docx_assignment(project['assignment_json'], include_answers=False)
            st.download_button("Download Assignment (.docx)", docx_student, "assignment.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        else:
            st.warning("No assignment available.")

    with tab_poster:
        if project['poster_json']:
            # Render simple HTML view of poster for student
            data = project['poster_json']
            st.markdown(f"## {data.get('poster_title')}")
            for sec in data.get('sections', []):
                st.markdown(f"#### {sec.get('heading')}")
                for b in sec.get('body_bullets', []):
                    st.markdown(f"- {b}")
            st.download_button("Download PDF", generate_pdf_poster(data), "poster.pdf", "application/pdf")
        else:
            st.warning("No reference material available.")

def main():
    init_db()
    
    # URL Query Param Handling for Students
    query_params = st.query_params
    token_param = query_params.get("token", None)

    with st.sidebar:
        st.title("ClassKit Studio 🍎")
        
        mode = st.radio("Mode", ["Teacher", "Student"])
        
        if mode == "Teacher":
            code = st.text_input("Access Code", type="password")
            if code == TEACHER_CODE:
                st.session_state['is_authenticated'] = True
                st.success("Verified Teacher")
            else:
                st.session_state['is_authenticated'] = False
                if code: st.error("Invalid Code")
        else:
            # Student Mode Instructions
            st.info("Paste your student link in the browser URL to view content.")

        st.divider()
        st.caption("Safety: This tool is for educational content generation. Do not use for academic dishonesty.")

    # --- ROUTING ---
    
    # 1. Student Route (Priority if token exists)
    if token_param:
        project = get_project_by_token(token_param, "student")
        if project:
            render_student_view(project)
        else:
            st.error("Invalid Student Token.")
        return

    # 2. Teacher Route
    if mode == "Teacher" and st.session_state.get('is_authenticated'):
        st.header("Teacher Dashboard")
        
        # New Project Form
        with st.expander("Create New Project", expanded=True):
            with st.form("new_proj"):
                title = st.text_input("Project Title")
                col1, col2 = st.columns(2)
                subj = col1.text_input("Subject")
                grade = col2.text_input("Grade Level")
                notes = st.text_area("Source Material / Topic Notes", height=100)
                safety = st.checkbox("I am generating materials for instruction, not cheating.")
                
                if st.form_submit_button("Create Project"):
                    if not safety:
                        st.error("Please confirm safety check.")
                    elif not title:
                        st.error("Title required.")
                    else:
                        new_id = str(uuid.uuid4())
                        t_token = str(uuid.uuid4())
                        s_token = str(uuid.uuid4())
                        proj_data = {
                            "id": new_id, "title": title, "subject": subj, "grade": grade,
                            "source_notes": notes, "teacher_token": t_token, "student_token": s_token
                        }
                        save_project(proj_data)
                        st.session_state['current_project_id'] = new_id
                        st.rerun()

        # Project Selector
        conn = sqlite3.connect(DB_FILE)
        conn.row_factory = sqlite3.Row
        projects = conn.cursor().execute("SELECT id, title, created_at FROM projects ORDER BY created_at DESC").fetchall()
        conn.close()

        if projects:
            p_opts = {p['id']: f"{p['title']} ({p['created_at'][:10]})" for p in projects}
            sel_id = st.selectbox("Select Project", options=list(p_opts.keys()), format_func=lambda x: p_opts[x], key="proj_sel")
            
            # Load selected or just created
            active_id = st.session_state.get('current_project_id', sel_id)
            project = load_project(active_id)
            
            if project:
                render_teacher_view(project)
        else:
            st.info("Create a project to get started.")

    elif mode == "Teacher":
        st.warning("Please enter access code in sidebar.")
    else:
        st.write("Welcome Student. Please use the link provided by your teacher.")

if __name__ == "__main__":
    main()